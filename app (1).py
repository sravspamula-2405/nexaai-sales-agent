# AI-Powered Sales Intelligence & Proposal Automation
# Gemini 3.5 Flash-Lite + URL Context / optional Google Search + local enterprise tools
import os, json, re, datetime, uuid
from pathlib import Path
from io import BytesIO
import pandas as pd
import streamlit as st

st.set_page_config(page_title="NexaAI Sales Intelligence", page_icon="🤖", layout="wide")
D = Path(__file__).parent / "data"

# -----------------------------
# Data repositories
# -----------------------------
K = pd.read_csv(D / "team_knowledge.csv")
J = pd.read_csv(D / "kanban.csv")
C = pd.read_csv(D / "crm.csv")
P = pd.read_csv(D / "products.csv")
CS = pd.read_csv(D / "case_studies.csv")
PR = pd.read_csv(D / "pricing.csv")
DM = pd.read_csv(D / "demos.csv")
IK = pd.read_csv(D / "industry_knowledge.csv")
COMP = pd.read_csv(D / "competitors.csv")
PART = pd.read_csv(D / "partners.csv")
SH = pd.read_csv(D / "sales_history.csv")
CTX = pd.read_csv(D / "industry_context.csv")

STOP = {"the","and","for","with","from","that","this","using","reduce","improve","high","customer","need","needs","want","looking","to","of","in","a","an","is","are","our","their","into","on","at","by","as","we"}

def tok(x):
    return set(re.findall(r"[a-zA-Z0-9]+", str(x).lower())) - STOP

def rank(df, q, n=None):
    x = df.copy()
    x["match"] = x.apply(lambda r: len(tok(" ".join(map(str, r.values))) & tok(q)), axis=1)
    x = x.sort_values("match", ascending=False)
    return x.head(n) if n else x

def df_records(df, n=6):
    if df is None or len(df) == 0:
        return []
    return df.head(n).drop(columns=["match"], errors="ignore").fillna("").to_dict(orient="records")

def safe(v, default="—"):
    t = "" if v is None else str(v).strip()
    return t if t else default

INDUSTRY_CONTEXTS = CTX["market"].tolist()

def context_row(market):
    m = CTX[CTX.market == market]
    return m.iloc[0].to_dict() if len(m) else {}

# Demo-friendly lead examples that change with the selected NexaAI sales context.
# These are fictional examples for demonstration purposes.
LEAD_DEFAULTS = {
    "fmcg": {"company":"FreshBasket Consumer Brands", "objective":"Improve demand forecasting and reduce inventory waste", "pain":"Demand uncertainty, stock-outs, overstocking, and limited supply-chain visibility", "website":""},
    "bank": {"company":"Apex Bank", "objective":"Improve fraud detection and automate risk monitoring", "pain":"Rising transaction fraud, manual investigation, and fragmented risk monitoring", "website":""},
    "health": {"company":"Apollo Care Network", "objective":"Reduce patient waiting time and improve hospital operations", "pain":"Long patient queues, inefficient resource allocation, and limited operational visibility", "website":""},
    "auto": {"company":"ABC Motors", "objective":"Reduce manufacturing downtime", "pain":"High equipment downtime, reactive maintenance, and limited asset visibility", "website":""},
    "it": {"company":"CloudWorks Digital", "objective":"Improve cloud operations and reduce IT infrastructure costs", "pain":"High cloud costs, fragmented infrastructure monitoring, and reactive IT operations", "website":""},
    "digital": {"company":"CloudWorks Digital", "objective":"Improve digital delivery predictability and data-driven decision making", "pain":"Fragmented delivery data, limited project visibility, and reactive decision making", "website":""},
    "data": {"company":"DataSphere Solutions", "objective":"Build a reliable enterprise data and AI operating model", "pain":"Fragmented data, inconsistent reporting, and difficulty scaling AI initiatives", "website":""},
}

def lead_defaults_for_market(market):
    m = str(market or "").lower()
    if "fmcg" in m or "consumer" in m: return LEAD_DEFAULTS["fmcg"].copy()
    if "bank" in m or "financial" in m: return LEAD_DEFAULTS["bank"].copy()
    if "health" in m or "pharma" in m: return LEAD_DEFAULTS["health"].copy()
    if "auto" in m or "mobility" in m: return LEAD_DEFAULTS["auto"].copy()
    if "it" in m: return LEAD_DEFAULTS["it"].copy()
    if "digital" in m: return LEAD_DEFAULTS["digital"].copy()
    if "data" in m or "ai" in m: return LEAD_DEFAULTS["data"].copy()
    ctx = context_row(market)
    return {"company":"New Customer Opportunity", "objective":f"Explore AI opportunities in {market}", "pain":safe(ctx.get("common_problems"), "Identify and address key operational challenges"), "website":""}

def filter_context(df, market):
    """Keep internal retrieval aligned to the selected go-to-market context."""
    if df is None or len(df) == 0 or not market:
        return df
    row = context_row(market)
    if not row:
        return df
    needles = [market.lower()] + [x.strip().lower() for x in str(row.get("sub_industries","")).split(";") if x.strip()]
    # Product/industry tables use a single industry field.
    if "industry" in df.columns:
        mask = df["industry"].astype(str).str.lower().apply(lambda x: any(n in x or x in n for n in needles))
        return df[mask]
    # Team/competitor/partner repositories use semicolon-separated industries.
    if "industries" in df.columns:
        mask = df["industries"].astype(str).str.lower().apply(lambda x: any(n in x or x in n for n in needles))
        return df[mask]
    return df


# -----------------------------
# Gemini helpers
# -----------------------------
def get_api_key():
    # Streamlit secrets first, then environment, then UI input.
    try:
        if "GEMINI_API_KEY" in st.secrets:
            return st.secrets["GEMINI_API_KEY"]
    except Exception:
        pass
    return os.getenv("GEMINI_API_KEY") or st.session_state.get("gemini_api_key", "")

@st.cache_resource(show_spinner=False)
def make_client(api_key):
    from google import genai
    return genai.Client(api_key=api_key)

AI_MODEL = "gemini-3.5-flash-lite"

# -----------------------------
# Tool definitions for the Orchestrator
# -----------------------------
TOOL_DEFS = [
    {"name":"search_crm","description":"Search customer and opportunity history in the CRM repository.","parameters":{"type":"object","properties":{"query":{"type":"string"}},"required":["query"]}},
    {"name":"search_products","description":"Find NexaAI products relevant to the customer's requirement, industry and pain point.","parameters":{"type":"object","properties":{"query":{"type":"string"}},"required":["query"]}},
    {"name":"search_industry","description":"Retrieve industry-specific problems, trends and AI opportunities.","parameters":{"type":"object","properties":{"query":{"type":"string"}},"required":["query"]}},
    {"name":"search_team_knowledge","description":"Find internal teams, skills and assets capable of delivering the proposed solution.","parameters":{"type":"object","properties":{"query":{"type":"string"}},"required":["query"]}},
    {"name":"search_projects","description":"Find relevant completed or active internal projects from the Kanban/project repository.","parameters":{"type":"object","properties":{"query":{"type":"string"}},"required":["query"]}},
    {"name":"search_case_studies","description":"Find similar customer case studies and business outcomes.","parameters":{"type":"object","properties":{"query":{"type":"string"}},"required":["query"]}},
    {"name":"search_pricing","description":"Find indicative commercial packages for selected products.","parameters":{"type":"object","properties":{"product_ids":{"type":"string"}},"required":["product_ids"]}},
    {"name":"search_demos","description":"Find the best available product demonstrations for selected products.","parameters":{"type":"object","properties":{"product_ids":{"type":"string"}},"required":["product_ids"]}},
    {"name":"search_competitors","description":"Find competitive positioning relevant to the lead.","parameters":{"type":"object","properties":{"query":{"type":"string"}},"required":["query"]}},
    {"name":"search_partners","description":"Find relevant integration or ecosystem partners.","parameters":{"type":"object","properties":{"query":{"type":"string"}},"required":["query"]}},
]

TOOL_MAP = {
    "search_crm": lambda q: rank(C, q, 4),
    "search_products": lambda q: rank(P, q, 5),
    "search_industry": lambda q: rank(IK, q, 3),
    "search_team_knowledge": lambda q: rank(K, q, 6),
    "search_projects": lambda q: rank(J, q, 6),
    "search_case_studies": lambda q: rank(CS, q, 4),
    "search_competitors": lambda q: rank(COMP, q, 3),
    "search_partners": lambda q: rank(PART, q, 4),
}

def execute_tool(name, args, market=None):
    try:
        if name in TOOL_MAP:
            raw = TOOL_MAP[name](safe(args.get("query"), ""))
            raw = filter_context(raw, market)
            return df_records(raw, 6)
        if name == "search_pricing":
            ids = [x.strip() for x in safe(args.get("product_ids"), "").split(",") if x.strip()]
            return df_records(PR[PR.product_id.astype(str).isin(ids)].drop_duplicates("product_id"), 6)
        if name == "search_demos":
            ids = [x.strip() for x in safe(args.get("product_ids"), "").split(",") if x.strip()]
            return df_records(DM[DM.product_id.astype(str).isin(ids)].drop_duplicates("product_id"), 6)
        return []
    except Exception as e:
        return {"error": str(e)}


def internal_agentic_orchestrator(client, lead, research, market):
    """LLM chooses which internal tools to call; application executes them; LLM synthesizes."""
    from google.genai import types

    declarations = [types.FunctionDeclaration(name=t["name"], description=t["description"], parameters=t["parameters"]) for t in TOOL_DEFS]
    # Force a real tool call on the first round so the demo visibly exercises the
    # agentic loop. Later rounds use AUTO so the model can decide when to stop.
    tool = types.Tool(function_declarations=declarations)
    config_auto = types.GenerateContentConfig(
        tools=[tool],
        tool_config=types.ToolConfig(
            function_calling_config=types.FunctionCallingConfig(mode="AUTO")
        ),
    )
    config_any = types.GenerateContentConfig(
        tools=[tool],
        tool_config=types.ToolConfig(
            function_calling_config=types.FunctionCallingConfig(mode="ANY")
        ),
    )

    system = f"""You are the Orchestrator Agent for NexaAI Technologies, an AI, Data & Digital Transformation company.\nThe selected go-to-market context is: {market}. Stay within this market context when selecting products, teams, projects and proof points.\nYour job is to build a customer-specific sales strategy from enterprise repositories.\nYou have tools for CRM, products, industry, teams, projects, case studies, pricing, demos, competitors and partners.\nDecide which tools are useful; do not call every tool automatically. Use the lead context and public research to target your searches.\nAfter retrieving evidence, produce a concise but specific sales strategy. Never invent product capabilities, prices, projects or case-study outcomes. If evidence is absent, say so.\nReturn these headings: Recommended Solution, Why It Fits, Customer-Specific Pitch, Proof Points, Commercial Approach, Demo Recommendation, Next Steps, Information Gaps."""
    user = f"""LEAD:\n{json.dumps(lead, ensure_ascii=False)}\n\nPUBLIC COMPANY RESEARCH:\n{research}\n\nUse the available enterprise tools to gather only the information needed to create a strong proposal."""
    contents = [types.Content(role="user", parts=[types.Part(text=system + "\n\n" + user)])]
    calls = []
    max_rounds = 5

    for _ in range(max_rounds):
        config = config_any if not calls else config_auto
        response = client.models.generate_content(model=AI_MODEL, contents=contents, config=config)
        candidate = response.candidates[0]
        contents.append(candidate.content)
        function_calls = [p.function_call for p in candidate.content.parts if getattr(p, "function_call", None)]
        if not function_calls:
            return response.text or "No AI synthesis was returned.", calls
        responses = []
        for fc in function_calls:
            args = dict(fc.args or {})
            result = execute_tool(fc.name, args, market)
            calls.append({"agent_tool": fc.name, "arguments": args, "records": result if isinstance(result, list) else []})
            # SDK compatibility: older google-genai releases do not accept the
            # optional function-call id. Newer releases do. Use it only when the
            # installed SDK supports it. GenerateContent expects the function
            # response in a user-role content message (not role="tool").
            try:
                import inspect
                accepts_id = "id" in inspect.signature(types.Part.from_function_response).parameters
            except Exception:
                accepts_id = False
            kwargs = {"name": fc.name, "response": {"result": result}}
            if accepts_id and getattr(fc, "id", None):
                kwargs["id"] = fc.id
            responses.append(types.Part.from_function_response(**kwargs))
        contents.append(types.Content(role="user", parts=responses))
    return "The agent reached its tool-call limit before producing a final synthesis.", calls


def research_company(client, company, industry, objective, pain, website=""):
    """Research the new lead. Prefer Gemini URL Context because it is available on the free tier;
    if no website is supplied, try Google Search and gracefully fall back to model analysis."""
    prompt = f"""Act as the Research Agent for a B2B sales opportunity.\nCompany: {company}\nIndustry: {industry}\nCustomer requirement: {objective}\nPain point: {pain}\nWebsite: {website or 'not supplied'}\n\nResearch and summarize only useful sales intelligence: what the company does, major operations, relevant business/technology priorities, likely operational needs connected to the stated requirement, and implications for an Industrial AI/Digital Transformation vendor. Clearly separate verified facts from sales hypotheses. Keep it concise and include source URLs when available."""
    try:
        if website.strip():
            # URL Context is free of charge and supported by Gemini 3.5 Flash-Lite.
            interaction = client.interactions.create(
                model=AI_MODEL,
                input=prompt + f"\n\nAnalyze the public content at this URL in depth: {website.strip()}",
                tools=[{"type": "url_context"}],
            )
            texts=[]
            for step in getattr(interaction, "steps", []) or []:
                if getattr(step, "type", "") == "model_output":
                    for block in getattr(step, "content", []) or []:
                        if getattr(block, "type", "") == "text":
                            texts.append(getattr(block, "text", ""))
            text="\n".join([x for x in texts if x]).strip()
            if text:
                return text, True
        # Optional broader search when available on the account.
        from google.genai import types
        response = client.models.generate_content(
            model=AI_MODEL,
            contents=prompt,
            config=types.GenerateContentConfig(tools=[types.Tool(google_search=types.ToolGoogleSearch())]),
        )
        return response.text or "No public research returned.", True
    except Exception as e:
        # Do not block the demo if a search-grounding entitlement is unavailable.
        try:
            response = client.models.generate_content(model=AI_MODEL, contents=prompt)
            text = response.text or ""
            if text:
                return text + "\n\n*Note: This fallback is model analysis, not live web-grounded research.*", True
        except Exception:
            pass
        return f"Public research unavailable: {e}", False


def fallback_ai_analysis(lead, research, internal):
    return ("## AI Sales Strategy\n\n"
            f"**Recommended Solution:** {safe(internal.get('products',''))}\n\n"
            "The live Gemini agent was not available, so the prototype used its structured repository matching as a fallback. "
            "Public research and internal evidence are shown separately so the salesperson can review them.")

# -----------------------------
# PPT renderer
# -----------------------------
def ppt(r):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.dml.color import RGBColor
    prs = Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    BG,DARK,MID,LIGHT,WHITE,ACCENT="F5F7FA","17202A","596675","DCE3EA","FFFFFF","2F5D8C"
    def bg(slide):
        s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height); s.fill.solid(); s.fill.fore_color.rgb=RGBColor.from_string(BG); s.line.fill.background()
        s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(.08)); s.fill.solid(); s.fill.fore_color.rgb=RGBColor.from_string(ACCENT); s.line.fill.background()
    def text(slide,x,y,w,h,value,size=16,bold=False,color=DARK,align=PP_ALIGN.LEFT):
        b=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=b.text_frame; tf.clear(); tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.TOP
        p=tf.paragraphs[0]; p.text=str(value); p.alignment=align; p.font.size=Pt(size); p.font.bold=bold; p.font.color.rgb=RGBColor.from_string(color); return b
    def title(slide,h,sub,page): text(slide,.65,.38,11.3,.52,h,27,True); text(slide,.67,.93,11,.36,sub,11,False,MID); text(slide,12.1,.42,.55,.25,page,9,False,MID,PP_ALIGN.RIGHT)
    def card(slide,x,y,w,h,hdr,body,hs=13,bs=12,accent=False):
        s=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h)); s.fill.solid(); s.fill.fore_color.rgb=RGBColor.from_string(WHITE); s.line.color.rgb=RGBColor.from_string(ACCENT if accent else LIGHT)
        text(slide,x+.20,y+.15,w-.40,.30,hdr,hs,True,ACCENT if accent else DARK); text(slide,x+.20,y+.55,w-.40,h-.68,body,bs,False,MID)
    def footer(slide): text(slide,.67,7.08,8,.18,"AI Sales Intelligence • Prototype • For discussion",8,False,MID)
    def slide(): s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s); return s
    def compact(v,n=180):
        t=safe(v,""); return t if len(t)<=n else t[:n-1].rstrip()+"…"
    def rv(row,col,default="—"):
        try: return safe(row[col],default) if col in row.index else default
        except: return default
    customer=safe(r.get("customer")); industry=safe(r.get("industry")); objective=safe(r.get("objective")); pain=safe(r.get("pain"))
    products=r.get("products",pd.DataFrame()); cases=r.get("cases",pd.DataFrame()); knowledge=r.get("knowledge",pd.DataFrame()); kanban=r.get("kanban",pd.DataFrame()); prices=r.get("prices",pd.DataFrame()); demos=r.get("demos",pd.DataFrame())
    ai=compact(r.get("ai_analysis",""),1500)
    # 1
    s=slide(); text(s,.85,1.05,11.4,.65,"AI-Powered Sales Proposal",34,True); text(s,.87,1.85,11,.75,customer,30,True,ACCENT); text(s,.89,2.70,10.5,.35,industry,15,False,MID); card(s,.85,3.75,11.6,1.55,"CUSTOMER OBJECTIVE",objective,12,20,True); text(s,.88,5.70,10.5,.45,"Research + organizational knowledge + AI orchestration",12,False,MID); footer(s)
    # 2
    s=slide(); title(s,"Understanding the Customer","Lead context provided by the salesperson","02"); card(s,.70,1.55,3.75,3.65,"BUSINESS OBJECTIVE",objective,13,18); card(s,4.78,1.55,3.75,3.65,"CURRENT PAIN POINT",pain,13,18); card(s,8.86,1.55,3.75,3.65,"INDUSTRY",industry,13,18); footer(s)
    # 3 AI strategy
    s=slide(); title(s,"AI-Generated Sales Strategy","Gemini analyzes public research and retrieved organizational evidence","03"); card(s,.70,1.40,7.45,4.55,"ORCHESTRATOR OUTPUT",compact(ai,1900),13,13,True); research=compact(r.get("research",""),1100); card(s,8.45,1.40,3.80,4.55,"PUBLIC COMPANY RESEARCH",research,13,12); footer(s)
    # 4 rec
    s=slide(); title(s,"Recommended Solution","Product intelligence ranked against the lead requirement","04"); rows=list(products.iterrows())[:3]; y=1.48
    for i,(_,p) in enumerate(rows): card(s,.70,y,7.75,1.25,f"{i+1}. {rv(p,'product','Solution')}",f"{compact(rv(p,'description',''),135)}\nIndicative price: {rv(p,'price','—')}  •  Relevance: {rv(p,'match','—')}",13,9,i==0); y+=1.43
    card(s,8.90,1.55,3.45,3.05,"SALES RECOMMENDATION","Use the AI strategy to position the highest-fit solution and alternatives around the customer's verified needs.",13,15,True); footer(s)
    # 5 capabilities
    s=slide(); title(s,"Why Us: Organizational Capabilities","Evidence retrieved from team knowledge","05"); y=1.42
    for _,k in list(knowledge.iterrows())[:5]: card(s,.70,y,5.85,.82,rv(k,'team','Team'),f"{compact(rv(k,'asset','Capability'),95)} • {rv(k,'status','Status')}",12,10); y+=.93
    card(s,7.00,1.50,5.25,3.35,"HOW THE AGENTS USE THIS","The Knowledge Agent identifies internal capabilities. The Orchestrator combines them with product, project and case-study evidence.",13,14,True); footer(s)
    # 6 project
    s=slide(); title(s,"Relevant Internal Project Evidence","What the organization has built or delivered","06"); y=1.42
    for _,k in list(kanban.iterrows())[:5]: card(s,.70,y,11.85,.82,compact(rv(k,'initiative','Initiative'),90),f"{rv(k,'team','Team')} • {rv(k,'status','Status')} • {compact(rv(k,'tags',''),70)}",12,10); y+=.93
    footer(s)
    # 7 case
    s=slide(); title(s,"Relevant Customer Experience","Proof point from the case-study repository","07")
    if len(cases):
        c=cases.iloc[0]; text(s,.85,1.48,7,.48,rv(c,'customer','Reference Customer'),24,True,ACCENT); text(s,.88,2.18,6.75,.25,"CHALLENGE",10,True,MID); text(s,.88,2.52,6.75,.85,compact(rv(c,'problem',''),180),15); text(s,.88,3.62,6.75,.25,"SOLUTION",10,True,MID); text(s,.88,3.96,6.75,.75,compact(rv(c,'solution',''),155),15,True); card(s,8.15,1.95,4.10,2.85,"RECORDED BUSINESS IMPACT",compact(rv(c,'impact',''),120),13,20,True)
    else: card(s,.85,1.65,11.4,3.0,"NO CASE STUDY FOUND","No relevant case study was returned.",15,18)
    footer(s)
    # 8 commercial
    s=slide(); title(s,"Indicative Commercial View","Pricing shown for prototype discussion only","08")
    for x,(_,pr) in zip([.80,6.85],PR[PR.product_id.isin(products.id)].drop_duplicates("product_id").head(2).iterrows() if len(products) else []):
        matches=products[products.id==pr.product_id]; pname=rv(matches.iloc[0],'product',rv(pr,'product_id','Solution')) if len(matches) else rv(pr,'product_id','Solution'); card(s,x,1.60,5.65,3.75,f"{pname} • {rv(pr,'package','Package')}",f"{rv(pr,'price','Indicative pricing')}\n\nMaximum indicative discount: {rv(pr,'max_discount','—')}\n\nValidate scope after discovery.",15,16,True)
    footer(s)
    # 9 next steps
    s=slide(); title(s,"Recommended Demo & Next Steps","Convert the proposal into an actionable customer conversation","09"); demo_text="\n".join(f"• {compact(rv(d,'demo','Focused demo'),145)}" for _,d in list(demos.head(3).iterrows())) or "• Focused solution demonstration"; card(s,.75,1.50,5.75,3.95,"RECOMMENDED DEMO",demo_text,14,16); card(s,6.85,1.50,5.75,3.95,"NEXT STEPS","1. Focused solution demonstration\n2. Technical discovery\n3. Validate data and integration needs\n4. Confirm commercials\n5. Human review and approval",14,16,True); footer(s)
    out=BytesIO(); prs.save(out); out.seek(0); return out

# -----------------------------
# Sidebar
# -----------------------------
source_map = {
    "CRM / Lead History": ("CRM / Lead History", C, ["Customer Agent", "Orchestrator"]),
    "Jira / Kanban": ("Jira / Kanban", J, ["Kanban Agent", "Orchestrator"]),
    "Team Knowledge": ("Team Knowledge", K, ["Knowledge Agent", "Orchestrator"]),
    "Product Knowledge": ("Product Knowledge", P, ["Product Agent", "Orchestrator"]),
    "Case Studies": ("Case Studies", CS, ["Case Study Agent", "Orchestrator"]),
    "Pricing Repository": ("Pricing Repository", PR, ["Pricing Agent", "Orchestrator"]),
    "Demo Repository": ("Demo Repository", DM, ["Demo Agent", "Orchestrator"]),
    "Industry Knowledge": ("Industry Knowledge", IK, ["Industry Agent", "Orchestrator"]),
    "Competitor Intelligence": ("Competitor Intelligence", COMP, ["Competitor Agent", "Orchestrator"]),
    "Partner Ecosystem": ("Partner Ecosystem", PART, ["Partner Agent", "Orchestrator"]),
    "Sales History": ("Sales History", SH, ["CRM Agent", "Knowledge Agent", "Orchestrator"]),
    "Industry Context": ("Industry Context", CTX, ["Research Agent", "Industry Agent", "Orchestrator"]),
}
if "selected_source" not in st.session_state: st.session_state.selected_source=None
with st.sidebar:
    st.title("🤖 NexaAI Sales Agent")
    st.caption("Gemini-powered agentic prototype")
    st.success(f"🎯 Context: {st.session_state.get('market_context','Not selected')}")
    if st.button("🔄 Change Industry Context", use_container_width=True):
        st.session_state.market_context = None
        st.rerun()
    st.divider()
    st.write("### 🔑 AI Agent Connection")
    key = st.text_input("Gemini API key", type="password", value=st.session_state.get("gemini_api_key", ""), help="Use your Gemini API key. It is kept in this session only unless you add it as a deployment secret.")
    st.session_state.gemini_api_key = key
    if key: st.success(f"API key entered • {AI_MODEL}")
    else: st.warning("No API key — deterministic prototype fallback")
    st.divider()
    st.write("### 🔗 Connected Enterprise Sources")
    for label in source_map:
        if st.button("🟢 " + label, key="src_"+label, use_container_width=True): st.session_state.selected_source=label
    st.divider(); st.write("### 🤖 Agent Team")
    for x in ["Research Agent","Customer Agent","Knowledge Agent","Kanban Agent","Product Agent","Case Study Agent","Pricing Agent","Demo Agent","Industry Agent","Competitor Agent","Partner Agent","Orchestrator Agent"]: st.write("• "+x)

# -----------------------------
# Opening: choose the go-to-market context
# -----------------------------
if "market_context" not in st.session_state:
    st.session_state.market_context = None

if not st.session_state.market_context:
    st.title("🏢 Configure Your Sales Intelligence Context")
    st.write("Before a lead arrives, define **what our company sells and which market we are targeting**. This industry context becomes the retrieval boundary and automatically loads an industry-relevant demo lead on the New Lead page.")
    st.divider()
    left, right = st.columns([1, 1.4])
    with left:
        market_choice = st.selectbox("Select our target industry / sales context", INDUSTRY_CONTEXTS, index=0)
        if st.button("🚀 Start Sales Workspace", type="primary", use_container_width=True):
            st.session_state.market_context = market_choice
            st.rerun()
    with right:
        ctx = context_row(market_choice)
        st.subheader(f"{market_choice} context")
        st.markdown(f"**Typical sub-industries:** {ctx.get('sub_industries','—')}")
        st.markdown(f"**Common customer problems:** {ctx.get('common_problems','—')}")
        st.markdown(f"**AI opportunities:** {ctx.get('ai_opportunities','—')}")
        st.markdown(f"**Typical buyers:** {ctx.get('typical_buyers','—')}")
        st.markdown(f"**Initial portfolio:** {ctx.get('portfolio','—')}")
    st.info("Prototype company context: **NexaAI Technologies — AI, Data & Digital Transformation solutions.** This selected market becomes the boundary for the downstream agent workflow.")
    st.stop()

market = st.session_state.market_context
st.title("AI-Powered Sales Intelligence & Proposal Automation")
st.info(f"🎯 Active sales context: **{market}**  |  New Lead → Public Company Research → Specialized Agents → Orchestrator → Customized Sales Pitch → PPT → Knowledge Enrichment")

if st.session_state.selected_source:
    label=st.session_state.selected_source; display_name,df,used_by=source_map[label]
    st.header(f"🔗 {display_name}")
    a,b,c=st.columns(3); a.metric("Connection","🟢 Connected"); b.metric("Records",len(df)); c.metric("Sync","Today")
    st.caption("Simulated enterprise repository used for the classroom prototype.")
    st.dataframe(df,hide_index=True,use_container_width=True)
    st.subheader("Agents using this source"); cols=st.columns(len(used_by))
    for col,agent in zip(cols,used_by): col.success("🤖 "+agent)
    st.markdown(f"**{display_name}** → relevant information is retrieved by the listed agents → the **Orchestrator Agent** synthesizes it into the lead-specific sales package.")
    if st.button("✖ Close source view"): st.session_state.selected_source=None; st.rerun()

tabs=st.tabs(["🚀 New Lead","🌐 AI Research","🧠 Agent Activity","📑 Sales Package","📊 Generate PPT","🔎 Audit Trail"])
if "r" not in st.session_state: st.session_state.r=None

with tabs[0]:
    st.header("New Lead")
    ctx=context_row(market)
    st.success(f"🏢 Selling into: **{market}** — {ctx.get('ai_opportunities','')} ")
    st.info("💡 The selected industry controls the initial customer example, likely problems, and the internal knowledge space searched by the agents.")

    # Load demo defaults only when the industry context changes. This prevents
    # Streamlit reruns from overwriting edits made by the salesperson.
    if st.session_state.get("_lead_defaults_market") != market:
        defaults = lead_defaults_for_market(market)
        st.session_state["lead_company"] = defaults["company"]
        st.session_state["lead_objective"] = defaults["objective"]
        st.session_state["lead_pain"] = defaults["pain"]
        st.session_state["lead_website"] = defaults.get("website", "")
        st.session_state["_lead_defaults_market"] = market

    c1,c2=st.columns(2)
    with c1:
        name=st.text_input("Company / Lead", key="lead_company")
        ind=st.selectbox("Customer Industry", [market], index=0, key="lead_industry", help="Aligned to the target industry context selected on the opening page.")
        website=st.text_input("Company website (optional)", key="lead_website")
    with c2:
        obj=st.text_input("Customer Requirement / Meeting Objective", key="lead_objective")
        pain=st.text_area("Customer Pain Points", key="lead_pain")
    st.caption(f"Industry-aware demo defaults loaded for **{market}**. The salesperson can edit any field before running the agents. The selected **{market}** context defines what NexaAI sells into this market; the Research Agent researches the company and the Orchestrator decides which specialized internal knowledge tools to use.")
    if st.button("✨ Run AI Sales Intelligence",type="primary",use_container_width=True):
        lead={"company":name,"industry":ind,"objective":obj,"pain":pain,"website":website,"target_sales_context":market}
        q=f"{market} {name} {ind} {obj} {pain}"
        api_key=get_api_key()
        research="Public company research unavailable because no Gemini API key was supplied."; research_ok=False; ai_text=""; calls=[]
        if api_key:
            try:
                client=make_client(api_key)
                with st.spinner("Research Agent is researching the company..."):
                    research,research_ok=research_company(client,name,ind,obj,pain,website)
                with st.spinner("Orchestrator is selecting internal agents and synthesizing the sales strategy..."):
                    ai_text,calls=internal_agentic_orchestrator(client,lead,research,market)
            except Exception as e:
                research=f"AI research failed: {e}"; ai_text=f"AI orchestration failed: {e}"; calls=[]
        # deterministic retrieval remains as a transparent data layer / fallback
        rr={"customer":name,"industry":ind,"objective":obj,"pain":pain,"website":website,
            "crm":filter_context(rank(C,q,3),market),"knowledge":filter_context(rank(K,q,10),market),"kanban":rank(J,q,6),"products":filter_context(rank(P,q,6),market).head(3),"cases":filter_context(rank(CS,q,6),market).head(3),
            "industry_knowledge":filter_context(rank(IK,q,4),market).head(2),"competitors":filter_context(rank(COMP,q,3),market).head(2),"partners":filter_context(rank(PART,q,4),market).head(3),"research":research,
            "research_ok":research_ok,"ai_analysis":ai_text,"tool_calls":calls}
        rr["prices"]=PR[PR.product_id.isin(rr["products"].id)].drop_duplicates("product_id")
        rr["demos"]=DM[DM.product_id.isin(rr["products"].id)].drop_duplicates("product_id")
        st.session_state.r=rr
        if api_key and ai_text and calls: st.success(f"Live AI agentic workflow completed • {len(calls)} internal tool calls")
        elif api_key: st.warning("Gemini returned a response, but no internal tool calls were recorded. The structured retrieval layer remains available.")
        else: st.warning("No Gemini key supplied. The prototype generated the structured fallback; add a key to activate live research + agent tool calling.")

with tabs[1]:
    r=st.session_state.r
    if r is None: st.info("Run a lead first.")
    else:
        st.header("🌐 AI Company Research")
        st.caption(f"Research Agent uses Gemini URL Context for the supplied company website. The result is interpreted against the selected **{market}** sales context; optional Google Search is attempted when available.")
        if r.get("research_ok"): st.success("Live public research completed")
        st.markdown(r.get("research","No research available."))

with tabs[2]:
    r=st.session_state.r
    if r is None: st.info("Run a lead first.")
    else:
        st.header("🧠 Agentic Orchestration")
        calls=r.get("tool_calls",[])
        if calls:
            for i,call in enumerate(calls,1):
                pretty=call["agent_tool"].replace("search_","").replace("_"," ").title()
                st.success(f"Step {i}: 🤖 **{pretty} Agent Tool** → retrieved {len(call.get('records',[]))} records")
                with st.expander("See what the agent requested"):
                    st.json(call.get("arguments",{}))
        else:
            st.info("No live tool calls were recorded. Add a Gemini API key and run the lead again to activate the LLM-driven tool-selection loop.")
        st.subheader("Orchestrator synthesis")
        st.markdown(r.get("ai_analysis", "No AI synthesis."))
        st.subheader("Retrieved internal evidence")
        st.caption(f"Retrieval is constrained by the active market context: **{market}**")
        st.dataframe(r["knowledge"],hide_index=True,use_container_width=True)
        st.dataframe(r["kanban"],hide_index=True,use_container_width=True)

with tabs[3]:
    r=st.session_state.r
    if r is None: st.info("Run a lead first.")
    else:
        st.header(f"Customized Sales Package — {r['customer']}")
        st.markdown(r.get("ai_analysis",""))
        st.subheader("Product Recommendations"); st.dataframe(r["products"][[c for c in ["product","description","price","match"] if c in r["products"].columns]],hide_index=True,use_container_width=True)
        st.subheader("Relevant Case Studies"); st.dataframe(r["cases"][[c for c in ["customer","problem","solution","impact"] if c in r["cases"].columns]],hide_index=True,use_container_width=True)
        st.subheader("Pricing"); st.dataframe(r["prices"],hide_index=True,use_container_width=True)
        st.subheader("Recommended Demo"); st.dataframe(r["demos"],hide_index=True,use_container_width=True)
        st.subheader("Industry Intelligence"); st.dataframe(r["industry_knowledge"],hide_index=True,use_container_width=True)
        st.subheader("Competitive Intelligence"); st.dataframe(r["competitors"],hide_index=True,use_container_width=True)
        st.subheader("Partner Ecosystem"); st.dataframe(r["partners"],hide_index=True,use_container_width=True)
        st.divider(); st.subheader("💾 Add Opportunity to Knowledge Base")
        st.caption("After human review, save this generated opportunity into the prototype CRM and sales-history repositories.")
        outcome=st.selectbox("Opportunity status",["Proposal Generated","Demo Scheduled","Negotiation","Won","Lost"],key="save_status")
        if st.button("💾 Save Approved Opportunity",type="primary",use_container_width=True):
            proposal_id="PROP-"+datetime.datetime.now().strftime("%Y%m%d-%H%M%S")
            primary=str(r["products"].iloc[0]["product"]) if len(r["products"]) else ""; secondary=str(r["products"].iloc[1]["product"]) if len(r["products"])>1 else ""; case_used=str(r["cases"].iloc[0]["customer"]) if len(r["cases"]) else ""; team_names=", ".join(dict.fromkeys(r["knowledge"]["team"].head(4).astype(str).tolist())) if len(r["knowledge"]) else ""; est=str(r["prices"].iloc[0]["price"]) if len(r["prices"]) else ""; today=datetime.date.today().isoformat()
            hist_row=pd.DataFrame([[proposal_id,r["customer"],r["industry"],r["objective"],primary,secondary,case_used,team_names,est,outcome,today,"AI-assisted proposal; human-reviewed before external sharing."]],columns=SH.columns); SH=pd.concat([SH,hist_row],ignore_index=True); SH.to_csv(D/"sales_history.csv",index=False)
            crm_id="CRM-"+datetime.datetime.now().strftime("%Y%m%d%H%M%S"); crm_row=pd.DataFrame([[crm_id,r["customer"],r["industry"],outcome,r["objective"],r["pain"],"50%",est,"Sales Intelligence"]],columns=C.columns); C=pd.concat([C,crm_row],ignore_index=True); C.to_csv(D/"crm.csv",index=False)
            source_map["CRM / Lead History"]=("CRM / Lead History",C,["Customer Agent","Orchestrator"]); source_map["Sales History"]=("Sales History",SH,["CRM Agent","Knowledge Agent","Orchestrator"])
            st.success(f"✅ Opportunity {proposal_id} saved to CRM and Sales History."); st.session_state.r["saved_proposal_id"]=proposal_id
        if r.get("saved_proposal_id"): st.info(f"Saved opportunity: **{r['saved_proposal_id']}** — open Sales History or CRM from the sidebar to see it.")
        st.subheader("Follow-up Email")
        primary_name=str(r["products"].iloc[0]["product"]) if len(r["products"]) else "the recommended solution"
        st.code(f"Subject: Follow-up – {r['objective']}\n\nHi {r['customer']} team,\n\nBased on our discussion, we recommend exploring {primary_name}. We have included relevant capabilities, proof points and indicative commercials. We would be happy to schedule a focused demonstration.\n\nRegards,\nSales Team")

with tabs[4]:
    r=st.session_state.r
    if r is None: st.info("Run a lead first.")
    else:
        st.header("Generate Customer Pitch Deck")
        st.write("Creates a fixed-layout customer deck using retrieved evidence plus the AI-generated strategy.")
        if st.button("📊 Generate PPT",type="primary"):
            try:
                f=ppt(r); st.download_button("⬇️ Download Customer Pitch Deck",f,file_name=f"{r['customer'].replace(' ','_')}_Sales_Pitch.pptx",mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            except ImportError: st.error("Run: pip install python-pptx")

with tabs[5]:
    r=st.session_state.r
    if r is None: st.info("Run a lead first.")
    else:
        st.header("🔎 Audit Trail")
        rows=[]
        for c in r.get("tool_calls",[]): rows.append([c["agent_tool"],"Local enterprise repository","LLM-selected retrieval",f"{len(c.get('records',[]))} records"])
        rows += [["Research Agent","Gemini URL Context / public website","Public company research","Completed" if r.get("research_ok") else "Unavailable"],["Orchestrator","All retrieved evidence","Sales strategy synthesis","Completed" if r.get("ai_analysis") else "Fallback"]]
        rows.append(["Human Review","Salesperson","Approval before external sharing","Required"])
        st.dataframe(pd.DataFrame(rows,columns=["Agent / Tool","Source","Contribution","Status"]),hide_index=True,use_container_width=True)
        st.warning("Production version: replace simulated CSV sources with authorized enterprise connectors, add access control, logging, data governance and persistent storage.")
